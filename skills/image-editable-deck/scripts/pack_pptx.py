"""
圖片簡報打包腳本

支援兩種模式：
- baked（預設）：圖裡已含文字，pptx 每頁一張 full-bleed 圖即可
- plate：圖為無文字底圖，依 YAML spec 疊加可編輯文字框
"""
import argparse
import glob
from pathlib import Path
import yaml
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def crop_to_ratio(src_path: str, target_w_in: float, target_h_in: float,
                  images_dir: Path = None) -> str:
    """依目標寬高比中央裁切，避免 python-pptx 強制拉伸變形。"""
    if not src_path:
        return src_path
    target_ratio = target_w_in / target_h_in
    img = Image.open(src_path)
    w, h = img.size
    current_ratio = w / h
    if abs(current_ratio - target_ratio) < 0.01:
        return src_path
    if current_ratio > target_ratio:
        new_w = int(h * target_ratio)
        left = (w - new_w) // 2
        img = img.crop((left, 0, left + new_w, h))
    else:
        new_h = int(w / target_ratio)
        top = (h - new_h) // 2
        img = img.crop((0, top, w, top + new_h))
    cache_dir = (images_dir or Path(src_path).parent) / "cropped"
    cache_dir.mkdir(parents=True, exist_ok=True)
    out = cache_dir / f"{Path(src_path).stem}__{target_w_in:.2f}x{target_h_in:.2f}.png"
    img.save(out)
    return str(out)


DEFAULT_PALETTE = {
    "bg": "#0D1B2A",
    "primary": "#00C6FF",
    "highlight": "#FFD700",
    "card": "#1E3A5F",
    "text": "#FFFFFF",
    "muted": "#A5B4CB",
}


def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def find_latest(images_dir: Path, prefix: str) -> str | None:
    cands = sorted(glob.glob(str(images_dir / f"{prefix}_*.png")))
    return cands[-1] if cands else None


def _resolve_color(name_or_hex: str, palette: dict) -> RGBColor:
    if not name_or_hex:
        return hex_to_rgb("#FFFFFF")
    if name_or_hex.startswith("#"):
        return hex_to_rgb(name_or_hex)
    hex_v = palette.get(name_or_hex, "#FFFFFF")
    return hex_to_rgb(hex_v)


def add_textbox(slide, block: dict, palette: dict, default_font: str, title_font: str | None = None, body_font: str | None = None):
    """依 block 規格在 slide 上加一個文字框。
    block keys: type, text, x, y, w, h, size, bold, color, align, anchor
    type: title / subtitle / body / badge / highlight
    """
    btype = block.get("type", "body")
    x = Inches(float(block["x"]))
    y = Inches(float(block["y"]))
    w = Inches(float(block["w"]))
    h = Inches(float(block["h"]))
    text = block.get("text", "")
    size = block.get("size")
    bold = block.get("bold")
    color_name = block.get("color")
    align_name = (block.get("align") or "left").lower()
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    anchor_name = (block.get("anchor") or "top").lower()

    # 類型預設值 — 依林長揚 #1 比例（55/34/21/13）規格化
    # 原比例在 16:9 投影片為海報級：title=55、subtitle=34、body=21、muted=13
    # 海報感適度放大至 title=72；其他按比例縮放
    defaults = {
        "title":     {"size": 72, "bold": True,  "color": "text"},
        "subtitle":  {"size": 34, "bold": True,  "color": "primary"},
        "body":      {"size": 21, "bold": False, "color": "text"},
        "badge":     {"size": 18, "bold": True,  "color": "bg"},
        "highlight": {"size": 26, "bold": True,  "color": "highlight"},
        "muted":     {"size": 14, "bold": False, "color": "muted"},
    }
    d = defaults.get(btype, defaults["body"])
    if d:
        size = size or d["size"]
        bold = d["bold"] if bold is None else bold
        color_name = color_name or d["color"]

    # 決定字型：block.font > 類型配對 > default_font
    TITLE_TYPES = {"title", "subtitle", "badge", "highlight"}
    font_name = block.get("font")
    if not font_name:
        if btype in TITLE_TYPES:
            font_name = title_font or default_font
        else:
            font_name = body_font or default_font

    # badge = 先畫圓角矩形底色，再放文字
    if btype == "badge":
        bg_color_name = block.get("bg") or "primary"
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _resolve_color(bg_color_name, palette)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Emu(30000)
        tf.margin_top = tf.margin_bottom = Emu(20000)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = text
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = _resolve_color(color_name, palette)
            run.font.name = font_name
        return

    # card = 矩形卡片底，無文字（僅背景）
    if btype == "card":
        card_color = block.get("bg") or "card"
        border_color = block.get("border")
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _resolve_color(card_color, palette)
        if border_color:
            shape.line.color.rgb = _resolve_color(border_color, palette)
            shape.line.width = Pt(block.get("border_width", 1.5))
        else:
            shape.line.fill.background()
        return

    # bar = 細橫條
    if btype == "bar":
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _resolve_color(color_name, palette)
        shape.line.fill.background()
        return

    # progress = 進度條（林長揚 #23：放進度條減輕觀眾壓力）
    # block 需有 current（目前頁）與 total（總頁數）
    if btype == "progress":
        current = int(block.get("current", 1))
        total = int(block.get("total", 10))
        track_color = block.get("track") or "card"
        fill_color = color_name or "primary"
        # 底條
        track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        track.fill.solid()
        track.fill.fore_color.rgb = _resolve_color(track_color, palette)
        track.line.fill.background()
        # 填充段
        ratio = max(0.0, min(1.0, current / total))
        fill_w = int(w.emu * ratio) if hasattr(w, "emu") else int(w * ratio)
        fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(fill_w), h)
        fill.fill.solid()
        fill.fill.fore_color.rgb = _resolve_color(fill_color, palette)
        fill.line.fill.background()
        return

    # 一般文字框
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor_map.get(anchor_name, MSO_ANCHOR.TOP)
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    # 林長揚 #5：行距為文字大小的 50–75%，預設 line_spacing=1.2（含字高 → 視覺行距約 60%）
    line_spacing = block.get("line_spacing", 1.2)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map.get(align_name, PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        p.text = line
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = _resolve_color(color_name, palette)
            run.font.name = font_name


def pack_baked(images_dir: Path, output: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    pngs = sorted(glob.glob(str(images_dir / "page_*.png")))
    if not pngs:
        raise SystemExit(f"錯誤：{images_dir} 找不到 page_*.png")

    by_page = {}
    for p in pngs:
        prefix = "_".join(Path(p).name.split("_")[:2])
        by_page[prefix] = p

    for prefix in sorted(by_page.keys()):
        png = by_page[prefix]
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(png, 0, 0, prs.slide_width, prs.slide_height)
        print(f"  [baked] {prefix}  <-  {Path(png).name}")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"[OK] {output.resolve()}  ({len(by_page)} 頁)")


def pack_plate(images_dir: Path, output: Path, spec_path: Path):
    spec = yaml.safe_load(spec_path.read_text(encoding="utf-8"))

    palette = {**DEFAULT_PALETTE, **(spec.get("style", {}).get("palette", {}))}
    style_cfg = spec.get("style", {})
    default_font = style_cfg.get("font", "Microsoft JhengHei")
    title_font = style_cfg.get("title_font") or default_font
    body_font = style_cfg.get("body_font") or default_font

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    pages = spec.get("pages", [])
    if not pages:
        raise SystemExit("spec.yaml 沒有 pages 欄位")

    for page_def in pages:
        slide = prs.slides.add_slide(blank)
        img_prefix = page_def.get("image")  # e.g. page_01
        img_path = find_latest(images_dir, img_prefix) if img_prefix else None

        # 全背景色底（保險用，底圖若透明或未對齊也看得乾淨）
        bg_name = page_def.get("bg") or "bg"
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = _resolve_color(bg_name, palette)
        bg_shape.line.fill.background()

        # 底圖（自動依目標 slot 比例裁切，避免變形）
        if img_path:
            img_w = float(page_def.get("img_w", 13.333))
            img_h = float(page_def.get("img_h", 7.5))
            ix = Inches(float(page_def.get("img_x", 0)))
            iy = Inches(float(page_def.get("img_y", 0)))
            # 若 page_def 明確指定 no_crop: true，則保留原圖
            if not page_def.get("no_crop"):
                img_path = crop_to_ratio(img_path, img_w, img_h, images_dir)
            slide.shapes.add_picture(img_path, ix, iy, Inches(img_w), Inches(img_h))

        # 文字層
        for block in page_def.get("blocks", []):
            add_textbox(slide, block, palette, default_font, title_font, body_font)

        print(f"  [plate] page {page_def.get('page')}  img={img_prefix}  blocks={len(page_def.get('blocks', []))}")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"[OK] {output.resolve()}  ({len(pages)} 頁)")


def main():
    p = argparse.ArgumentParser(description="圖片簡報打包 pptx")
    p.add_argument("--images-dir", default="slides/images", help="圖片目錄")
    p.add_argument("--output", default="slides/output.pptx", help="輸出 pptx")
    p.add_argument("--mode", choices=["baked", "plate"], default="baked",
                   help="baked=圖內含文字；plate=底圖+可編輯文字框")
    p.add_argument("--spec", default=None, help="plate 模式的 YAML 規格檔")
    args = p.parse_args()

    images_dir = Path(args.images_dir)
    output = Path(args.output)

    if args.mode == "baked":
        pack_baked(images_dir, output)
    else:
        if not args.spec:
            raise SystemExit("plate 模式需要 --spec <spec.yaml>")
        pack_plate(images_dir, output, Path(args.spec))


if __name__ == "__main__":
    main()
